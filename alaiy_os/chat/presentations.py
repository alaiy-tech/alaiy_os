"""Turning a slide outline the model has written into a .pptx deck.

The presentation sibling of `exports.py`: a normalised structure in, bytes
out, no Frappe documents — so, like that module, it can be exercised without
a site. The Frappe side (the `File` row, the permission, the download chip)
lives in `artifacts.py`, in `create_presentation()`.

## Why a separate tool, not a fourth `create_download` format

`create_download`'s schema is a table: columns and rows. A presentation is
not a table — it's an ordered sequence of slides, each with its own shape.
Forcing that into columns/rows would mean encoding slide structure as magic
column names, which is exactly the kind of implicit contract this codebase
avoids everywhere else. One tool, one schema that matches what it actually
produces — the same reasoning `create_download`'s own docstring gives for
being one tool with a format enum rather than three.

## Slide types

    title:   heading, subtitle (optional) — meant for the first slide.
    bullets: heading, bullets (a list of strings).
    table:   heading, columns, rows — the same columns/rows shape
             `create_download` uses, so the model reuses what it already
             knows rather than learning a second way to describe a table.
    chart:   heading, chart_type (bar/line/pie), labels, series, unit — the
             same shape as the inline ```alaiy-chart``` fence convention (see
             `runner.CHART_PROMPT`), field-for-field, except "type" is
             renamed to "chart_type" here because "type" already names the
             *slide* kind ("chart") at this level. A native, editable
             PowerPoint chart, not a picture of one — see `_add_chart_slide`.

Table and chart slides are capped much smaller than their chat-bubble/
create_download equivalents (`MAX_TABLE_ROWS` vs. `exports.MAX_ROWS`,
`MAX_CHART_LABELS` vs. the chat fence's 60): a slide is read at a glance from
across a room, not scrolled or zoomed — past a dozen rows or categories it
stops being a slide and starts being a spreadsheet or chat chart someone
pasted in, which is what `create_download` and the inline fence are for.
Every throw below says so.

A slide with an unrecognised type is refused outright, not skipped: silently
dropping content the user asked for is a worse failure than a visible error
the model can read and correct.

## Visual design

Deliberately not python-pptx's default white-background, black-Calibri
look. Every slide gets a slim accent bar and consistent heading colour/size
(`_ensure_styles`... no, wrong module — see `_ACCENT`, `_add_accent_bar`),
table slides get a coloured header row and banded body rows, and charts get
data labels, a formatted axis, and the same accent colour as everything
else, so a generated deck reads as one deliberately designed document
instead of a template with the model's words dropped in.
"""

import io
import os
import re

import frappe

MAX_SLIDES = 20
MAX_BULLETS_PER_SLIDE = 8
MAX_BULLET_CHARS = 200
MAX_HEADING_CHARS = 100

# A table slide is read at a glance, not scrolled -- see the module docstring.
MAX_TABLE_ROWS = 12
MAX_TABLE_COLUMNS = 6

# Same reasoning as the table caps -- see the module docstring.
MAX_CHART_SERIES = 4
MAX_CHART_LABELS = 20
MAX_PIE_LABELS = 8
MAX_SERIES_NAME_CHARS = 40

# Checked after generation, the same way exports.MAX_BYTES is: embedded
# fonts/theme data make a .pptx's size a function of python-pptx's template,
# not just of how much text went in, so it's only knowable once written.
MAX_BYTES = 5 * 1024 * 1024

MAX_FILE_NAME_CHARS = 80

# Anything that is not plainly part of a file name. Mirrors
# exports._UNSAFE_NAME, duplicated rather than imported: that one is private
# to exports.py, and the two modules should stay free to diverge (e.g. if a
# future format needs a character this one doesn't).
_UNSAFE_NAME = re.compile(r"[^A-Za-z0-9 ._+-]+")

SLIDE_TYPES = ("title", "bullets", "table", "chart")
CHART_TYPES = ("bar", "line", "pie")
CHART_UNITS = ("number", "currency", "percent")


def safe_file_name(name):
	"""A file name from a language model, made safe -- no extension.

	Unlike `exports.safe_file_name`, which looks up the extension in
	`exports.EXTENSIONS`: pptx deliberately isn't one of those (see this
	module's own docstring), so the caller appends `.pptx` itself.
	"""
	base = os.path.basename(str(name or "").strip()) or "presentation"
	base = os.path.splitext(base)[0]
	base = _UNSAFE_NAME.sub("-", base).strip(" .-") or "presentation"
	return base[:MAX_FILE_NAME_CHARS]


def _text(value, limit):
	"""One value as a capped string. `None` is empty, not "None"."""
	if value is None:
		return ""
	text = value if isinstance(value, str) else str(value)
	if len(text) > limit:
		return text[:limit] + "…"
	return text


def _number(value, context):
	"""One chart value as a float, or `None` for a genuinely missing point.

	`None` round-trips (a gap in the series, exactly like the inline
	```alaiy-chart``` fence's "use null for a value you do not have"); anything
	else that cannot become a number is refused rather than silently zeroed,
	the same way a bad row shape is refused in `normalise`.
	"""
	if value is None:
		return None
	if isinstance(value, bool):
		frappe.throw(frappe._('{0}: chart values must be numbers, not "true"/"false".').format(context))
	try:
		return float(value)
	except (TypeError, ValueError):
		frappe.throw(frappe._('{0}: "{1}" is not a number.').format(context, value))


def normalise(slides):
	"""`slides`, validated and every cap applied. Throws when one is hit.

	Mirrors `exports.normalise`: every throw is written for the model to read
	back as a `tool_result` and act on, since that is its only route to
	recovering from a bad call.
	"""
	if not slides:
		frappe.throw(frappe._("A presentation needs at least one slide."))
	if len(slides) > MAX_SLIDES:
		frappe.throw(
			frappe._("{0} slides is more than the {1} a presentation can hold. Cut it down.").format(
				len(slides), MAX_SLIDES
			)
		)

	out = []
	for index, slide in enumerate(slides):
		position = index + 1
		if not isinstance(slide, dict):
			frappe.throw(
				frappe._("Slide {0} must be an object with at least a type and a heading.").format(position)
			)

		kind = slide.get("type")
		if kind not in SLIDE_TYPES:
			frappe.throw(
				frappe._('Slide {0} has type "{1}", which is not one of: {2}.').format(
					position, kind, ", ".join(SLIDE_TYPES)
				)
			)

		heading = _text(slide.get("heading"), MAX_HEADING_CHARS)
		if not heading:
			frappe.throw(frappe._("Slide {0} needs a heading.").format(position))

		if kind == "title":
			out.append(
				{
					"type": "title",
					"heading": heading,
					"subtitle": _text(slide.get("subtitle"), MAX_HEADING_CHARS),
				}
			)
			continue

		if kind == "bullets":
			bullets = [b for b in (_text(b, MAX_BULLET_CHARS) for b in (slide.get("bullets") or [])) if b]
			if not bullets:
				frappe.throw(
					frappe._('Slide {0} ("{1}") is type "bullets" but has no bullets.').format(position, heading)
				)
			if len(bullets) > MAX_BULLETS_PER_SLIDE:
				frappe.throw(
					frappe._(
						'Slide {0} ("{1}") has {2} bullets, more than the {3} that fit on one '
						"slide. Split it across slides."
					).format(position, heading, len(bullets), MAX_BULLETS_PER_SLIDE)
				)
			out.append({"type": "bullets", "heading": heading, "bullets": bullets})
			continue

		if kind == "table":
			columns = [_text(c, MAX_HEADING_CHARS) for c in (slide.get("columns") or [])]
			if not columns:
				frappe.throw(
					frappe._('Slide {0} ("{1}") is type "table" but has no columns.').format(position, heading)
				)
			if len(columns) > MAX_TABLE_COLUMNS:
				frappe.throw(
					frappe._(
						'Slide {0} ("{1}") has {2} columns, more than the {3} that fit on one '
						"slide. Drop some, or use create_download for the full table."
					).format(position, heading, len(columns), MAX_TABLE_COLUMNS)
				)

			rows = []
			for row in slide.get("rows") or []:
				if not isinstance(row, (list, tuple)):
					frappe.throw(
						frappe._('Slide {0} ("{1}"): every row must be a list of values, one per column.').format(
							position, heading
						)
					)
				values = [_text(v, MAX_HEADING_CHARS) for v in row[: len(columns)]]
				values += [""] * (len(columns) - len(values))
				rows.append(values)
			if not rows:
				frappe.throw(
					frappe._('Slide {0} ("{1}") is type "table" but has no rows.').format(position, heading)
				)
			if len(rows) > MAX_TABLE_ROWS:
				frappe.throw(
					frappe._(
						'Slide {0} ("{1}") has {2} rows, more than the {3} that fit on one slide. '
						"Summarise, or use create_download for the full table."
					).format(position, heading, len(rows), MAX_TABLE_ROWS)
				)

			out.append({"type": "table", "heading": heading, "columns": columns, "rows": rows})
			continue

		# chart
		context = f'Slide {position} ("{heading}")'
		chart_type = slide.get("chart_type") or "bar"
		if chart_type not in CHART_TYPES:
			frappe.throw(
				frappe._('{0}: chart_type "{1}" is not one of: {2}.').format(
					context, chart_type, ", ".join(CHART_TYPES)
				)
			)

		unit = slide.get("unit") or "number"
		if unit not in CHART_UNITS:
			frappe.throw(frappe._('{0}: unit "{1}" is not one of: {2}.').format(context, unit, ", ".join(CHART_UNITS)))

		labels = [_text(v, MAX_HEADING_CHARS) for v in (slide.get("labels") or [])]
		if not labels:
			frappe.throw(frappe._("{0} is type \"chart\" but has no labels.").format(context))
		label_limit = MAX_PIE_LABELS if chart_type == "pie" else MAX_CHART_LABELS
		if len(labels) > label_limit:
			frappe.throw(
				frappe._(
					"{0} has {1} labels, more than the {2} a {3} chart can hold on one slide. "
					"Chart the top few and mention the rest in a bullets slide instead."
				).format(context, len(labels), label_limit, chart_type)
			)

		raw_series = slide.get("series") or []
		if not raw_series:
			frappe.throw(frappe._("{0} is type \"chart\" but has no series.").format(context))
		if len(raw_series) > MAX_CHART_SERIES:
			frappe.throw(
				frappe._("{0} has {1} series, more than the {2} one chart can hold.").format(
					context, len(raw_series), MAX_CHART_SERIES
				)
			)
		if chart_type == "pie" and len(raw_series) != 1:
			frappe.throw(
				frappe._('{0}: a pie chart takes exactly one series (it shows shares of one total).').format(context)
			)

		series = []
		for s_index, s in enumerate(raw_series):
			if not isinstance(s, dict):
				frappe.throw(frappe._("{0}: every series must be an object with a name and points.").format(context))
			name = _text(s.get("name"), MAX_SERIES_NAME_CHARS) or f"Series {s_index + 1}"
			points = s.get("points") or []
			if len(points) != len(labels):
				frappe.throw(
					frappe._(
						'{0}: series "{1}" has {2} points but there are {3} labels -- every series '
						"needs exactly one point per label, in the same order."
					).format(context, name, len(points), len(labels))
				)
			series.append(
				{"name": name, "points": [_number(v, f'{context}, series "{name}"') for v in points]}
			)

		out.append(
			{
				"type": "chart",
				"heading": heading,
				"chart_type": chart_type,
				"unit": unit,
				"labels": labels,
				"series": series,
			}
		)

	return out


# ── Visual design ──────────────────────────────────────────────────────────
# One accent colour, used everywhere something needs to look deliberate
# rather than default: heading text, the top bar on every content slide, a
# table's header row, a chart's series and axis. Consistency is what reads as
# "designed" from a template with no real theme behind it -- see the module
# docstring.
def _accent():
	from pptx.dml.color import RGBColor

	return RGBColor(0x1F, 0x4E, 0x79)


def _accent_light():
	from pptx.dml.color import RGBColor

	return RGBColor(0xEA, 0xF1, 0xFA)


def _ink():
	from pptx.dml.color import RGBColor

	return RGBColor(0x22, 0x22, 0x22)


def _add_accent_bar(slide, deck):
	"""A slim coloured strip along the top of a content slide.

	The one repeated element that makes a sequence of otherwise-plain slides
	read as one deck rather than a template -- see the module docstring.
	Skipped on the title slide, which gets its own full-width treatment
	instead (see `_add_title_slide`).
	"""
	from pptx.enum.shapes import MSO_SHAPE
	from pptx.util import Emu

	bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), deck.slide_width, Emu(91440))
	bar.fill.solid()
	bar.fill.fore_color.rgb = _accent()
	bar.line.fill.background()
	bar.shadow.inherit = False
	return bar


def _add_title_slide(deck, slide):
	from pptx.enum.shapes import MSO_SHAPE
	from pptx.util import Emu, Inches, Pt

	s = deck.slides.add_slide(deck.slide_layouts[0])
	s.shapes.title.text = slide["heading"]
	title_frame = s.shapes.title.text_frame
	title_frame.paragraphs[0].font.size = Pt(40)
	title_frame.paragraphs[0].font.bold = True
	title_frame.paragraphs[0].font.color.rgb = _accent()

	if slide["subtitle"] and len(s.placeholders) > 1:
		s.placeholders[1].text = slide["subtitle"]
		s.placeholders[1].text_frame.paragraphs[0].font.color.rgb = _ink()

	# A full-width accent band along the bottom -- the title slide's own
	# design flourish, distinct from every other slide's top bar, so the
	# first slide reads as a cover rather than just another content slide.
	band = s.shapes.add_shape(
		MSO_SHAPE.RECTANGLE, Emu(0), deck.slide_height - Inches(0.35), deck.slide_width, Inches(0.35)
	)
	band.fill.solid()
	band.fill.fore_color.rgb = _accent()
	band.line.fill.background()
	band.shadow.inherit = False
	return s


def _add_bullets_slide(deck, slide):
	from pptx.util import Pt

	s = deck.slides.add_slide(deck.slide_layouts[1])
	s.shapes.title.text = slide["heading"]
	s.shapes.title.text_frame.paragraphs[0].font.color.rgb = _accent()
	s.shapes.title.text_frame.paragraphs[0].font.bold = True

	body = s.placeholders[1].text_frame
	body.clear()
	for i, bullet in enumerate(slide["bullets"]):
		paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
		paragraph.text = bullet
		paragraph.font.size = Pt(20)
		paragraph.font.color.rgb = _ink()

	_add_accent_bar(s, deck)
	return s


def _add_heading_box(s, deck, heading):
	"""The heading textbox shared by table and chart slides -- neither uses
	a layout's own title placeholder, since both need the space below it
	sized in absolute inches (`add_table`/`add_chart` take explicit
	coordinates regardless of any placeholder)."""
	from pptx.util import Inches, Pt

	box = s.shapes.add_textbox(Inches(0.5), Inches(0.35), deck.slide_width - Inches(1), Inches(0.8))
	frame = box.text_frame
	frame.text = heading
	frame.paragraphs[0].font.size = Pt(28)
	frame.paragraphs[0].font.bold = True
	frame.paragraphs[0].font.color.rgb = _accent()
	return box


def _add_table_slide(deck, slide):
	from pptx.util import Inches, Pt

	s = deck.slides.add_slide(deck.slide_layouts[6])
	_add_accent_bar(s, deck)
	_add_heading_box(s, deck, slide["heading"])

	row_count, col_count = len(slide["rows"]) + 1, len(slide["columns"])
	table_shape = s.shapes.add_table(
		row_count, col_count, Inches(0.5), Inches(1.4), deck.slide_width - Inches(1), Inches(0.5 * row_count)
	)
	table = table_shape.table

	for col_index, heading in enumerate(slide["columns"]):
		cell = table.cell(0, col_index)
		cell.text = heading
		cell.fill.solid()
		cell.fill.fore_color.rgb = _accent()
		para = cell.text_frame.paragraphs[0]
		para.font.bold = True
		para.font.size = Pt(14)
		from pptx.dml.color import RGBColor

		para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

	# Banded body rows -- a plain white grid is the single easiest thing to
	# misread across a room; alternating a pale tint of the accent colour is
	# what every real slide-table template does, for that reason.
	for row_index, row in enumerate(slide["rows"], start=1):
		banded = row_index % 2 == 0
		for col_index, value in enumerate(row):
			cell = table.cell(row_index, col_index)
			cell.text = value
			cell.text_frame.paragraphs[0].font.size = Pt(13)
			cell.text_frame.paragraphs[0].font.color.rgb = _ink()
			if banded:
				cell.fill.solid()
				cell.fill.fore_color.rgb = _accent_light()
			else:
				cell.fill.solid()
				from pptx.dml.color import RGBColor

				cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

	return s


#: unit -> the number format string applied to a chart's axis and data labels.
_NUMBER_FORMATS = {"number": "#,##0", "currency": "$#,##0", "percent": "0%"}
#: unit -> the (type, chart_type) pair in `pptx.enum.chart.XL_CHART_TYPE`.


def _add_chart_slide(deck, slide):
	from pptx.chart.data import CategoryChartData
	from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
	from pptx.util import Inches, Pt

	s = deck.slides.add_slide(deck.slide_layouts[6])
	_add_accent_bar(s, deck)
	_add_heading_box(s, deck, slide["heading"])

	chart_data = CategoryChartData()
	chart_data.categories = slide["labels"]
	for series in slide["series"]:
		# python-pptx accepts None entries directly -- they render as a gap,
		# exactly the "use null for a value you do not have" contract the
		# inline ```alaiy-chart``` fence already promises the model.
		chart_data.add_series(series["name"], series["points"])

	xl_type = {
		"bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
		"line": XL_CHART_TYPE.LINE_MARKERS,
		"pie": XL_CHART_TYPE.PIE,
	}[slide["chart_type"]]

	graphic_frame = s.shapes.add_chart(
		xl_type, Inches(0.6), Inches(1.5), deck.slide_width - Inches(1.2), deck.slide_height - Inches(2.1), chart_data
	)
	chart = graphic_frame.chart
	chart.has_title = False  # the heading textbox above already says this

	number_format = _NUMBER_FORMATS[slide["unit"]]
	plot = chart.plots[0]
	plot.has_data_labels = True
	plot.data_labels.number_format = number_format
	plot.data_labels.number_format_is_linked = False
	plot.data_labels.font.size = Pt(11)

	if slide["chart_type"] == "pie":
		# A pie's "series" is one ring of shares -- colour each slice, not
		# each series, and a legend genuinely helps since slice labels alone
		# get crowded past a handful of categories.
		points = plot.series[0].points
		palette = _pie_palette(len(points))
		for point, color in zip(points, palette):
			point.format.fill.solid()
			point.format.fill.fore_color.rgb = color
		chart.has_legend = True
		chart.legend.position = XL_LEGEND_POSITION.RIGHT
		chart.legend.include_in_layout = False
		chart.legend.font.size = Pt(11)
	else:
		for series in chart.series:
			series.format.fill.solid()
			series.format.fill.fore_color.rgb = _accent()
		if len(chart.series) > 1:
			chart.has_legend = True
			chart.legend.position = XL_LEGEND_POSITION.BOTTOM
			chart.legend.include_in_layout = False
			chart.legend.font.size = Pt(11)
		value_axis = chart.value_axis
		value_axis.tick_labels.number_format = number_format
		value_axis.tick_labels.number_format_is_linked = False
		value_axis.tick_labels.font.size = Pt(11)
		chart.category_axis.tick_labels.font.size = Pt(11)

	return s


def _pie_palette(count):
	"""`count` shades of the accent colour, lightest to darkest, so a pie's
	slices are visibly distinct without reaching for unrelated hues that
	would clash with the rest of the deck."""
	from pptx.dml.color import RGBColor

	base = (0x1F, 0x4E, 0x79)
	colors = []
	for i in range(count):
		# Blend toward white as i grows, capped so the last slice is still
		# clearly the accent colour rather than washed out to nothing.
		t = 0.75 * (i / max(count - 1, 1))
		colors.append(RGBColor(*(int(c + (255 - c) * t) for c in base)))
	return colors


def write_pptx(slides, deck_title=None):
	"""Bytes of one .pptx file. `slides` must have already been through `normalise`.

	`deck_title` is currently unused by the writer itself — python-pptx has no
	document-properties call cheap enough to be worth a second code path here —
	but is accepted for the same shape `exports.write` takes, and is where core
	properties (title/author) would go if a reader ever asked for them.
	"""
	from pptx import Presentation

	deck = Presentation()
	# 16:9, not python-pptx's 4:3 default -- every deck viewer opens widescreen
	# now, and a 4:3 export reads as a decade out of date.
	from pptx.util import Inches

	deck.slide_width = Inches(13.333)
	deck.slide_height = Inches(7.5)

	builders = {
		"title": _add_title_slide,
		"bullets": _add_bullets_slide,
		"table": _add_table_slide,
		"chart": _add_chart_slide,
	}
	for slide in slides:
		builders[slide["type"]](deck, slide)

	buffer = io.BytesIO()
	deck.save(buffer)
	return buffer.getvalue()
